#!/usr/bin/env python3
"""Generate a graphics-first lab-meeting PPTX from balanced ML benchmark and gallery artifacts."""

from __future__ import annotations

import argparse
import json
from pathlib import Path
from typing import Any

import matplotlib.pyplot as plt
import numpy as np
from matplotlib import patches
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


BLACK = RGBColor(0, 0, 0)
WHITE = RGBColor(255, 255, 255)
TEXT = RGBColor(20, 20, 20)
GRAY = RGBColor(230, 230, 230)


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description="Generate a lab-meeting executive assessment PPTX.")
    parser.add_argument(
        "--benchmark-root",
        type=Path,
        default=Path("reports/ml/benchmarks/data_march2026_balanced"),
        help="Balanced benchmark root directory.",
    )
    parser.add_argument(
        "--gallery-root",
        type=Path,
        default=Path("reports/ml/diagnostic_gallery/example"),
        help="Diagnostic gallery output directory.",
    )
    parser.add_argument(
        "--dataset-manifest",
        type=Path,
        default=Path("reports/ml/datasets/data_march2026_balanced/manifest.json"),
        help="Balanced dataset manifest path.",
    )
    parser.add_argument(
        "--output-dir",
        type=Path,
        default=Path("reports/ml/presentations"),
        help="Output directory for PPTX, manifest, and generated figures.",
    )
    parser.add_argument(
        "--basename",
        type=str,
        default="data-march-2026-balanced-executive-assessment",
        help="Output basename without extension.",
    )
    parser.add_argument(
        "--deck-title",
        type=str,
        default="Balanced Ni-Cu-Al ML Assessment",
        help="Deck title.",
    )
    return parser.parse_args()


def _load_json(path: Path) -> dict[str, Any]:
    return json.loads(path.read_text(encoding="utf-8"))


def _setup_matplotlib() -> None:
    plt.rcParams.update(
        {
            "font.family": "Arial",
            "font.size": 18,
            "axes.titlesize": 20,
            "axes.labelsize": 18,
            "xtick.labelsize": 16,
            "ytick.labelsize": 16,
            "figure.facecolor": "white",
            "axes.facecolor": "white",
        }
    )


def _savefig(fig: plt.Figure, path: Path) -> Path:
    path.parent.mkdir(parents=True, exist_ok=True)
    fig.savefig(path, dpi=200, bbox_inches="tight", facecolor="white")
    plt.close(fig)
    return path


def _mean(values: list[float]) -> float:
    return float(sum(values) / max(1, len(values)))


def _gallery_stats(gallery_manifest: dict[str, Any]) -> dict[str, dict[str, Any]]:
    by_source: dict[str, list[dict[str, Any]]] = {}
    for record in gallery_manifest["records"]:
        by_source.setdefault(str(record["source_label"]), []).append(record)

    out: dict[str, dict[str, Any]] = {}
    for source, rows in by_source.items():
        counts: dict[str, int] = {}
        for row in rows:
            phase = str(row["predicted_phase"])
            counts[phase] = counts.get(phase, 0) + 1
        out[source] = {
            "counts": counts,
            "mean_confidence": _mean([float(row["confidence"]) for row in rows]),
            "mean_margin": _mean([float(row["margin"]) for row in rows]),
        }
    return out


def _figure_workflow(path: Path) -> Path:
    fig, ax = plt.subplots(figsize=(14, 4.4))
    ax.set_axis_off()

    boxes = [
        ("Single-phase\n.oh5 scans", 0.03, "#d9eaf7"),
        ("Quality filter\nCI > 0.2, Fit < 1.5", 0.20, "#fce6c8"),
        ("Phase balancing\n1751 per phase", 0.40, "#e4f1d5"),
        ("Stratified split\n4200 / 525 / 528", 0.58, "#f4d6d6"),
        ("5-model benchmark\nsuite", 0.75, "#d9dbf7"),
        ("Diagnostic gallery\nreference + blind", 0.89, "#d8f0ea"),
    ]

    y = 0.34
    h = 0.34
    w = 0.13
    for idx, (label, x, color) in enumerate(boxes):
        rect = patches.FancyBboxPatch(
            (x, y),
            w,
            h,
            boxstyle="round,pad=0.02,rounding_size=0.02",
            linewidth=2,
            edgecolor="black",
            facecolor=color,
        )
        ax.add_patch(rect)
        ax.text(x + w / 2, y + h / 2, label, ha="center", va="center", fontsize=18, weight="bold")
        if idx < len(boxes) - 1:
            ax.annotate(
                "",
                xy=(boxes[idx + 1][1] - 0.01, y + h / 2),
                xytext=(x + w + 0.005, y + h / 2),
                arrowprops=dict(arrowstyle="->", lw=2.5, color="black"),
            )

    ax.text(0.03, 0.88, "Balanced March 2026 ML pipeline", fontsize=22, weight="bold")
    ax.text(0.03, 0.16, "Training set: 1400 Al + 1400 Cu + 1400 Ni | Selected gallery model: simple_cnn_w32", fontsize=18)
    return _savefig(fig, path)


def _figure_dataset(dataset_manifest: dict[str, Any], path: Path) -> Path:
    fig = plt.figure(figsize=(13.5, 7.2))
    gs = fig.add_gridspec(2, 2, height_ratios=[1.1, 0.95], width_ratios=[1.1, 1.0], hspace=0.35, wspace=0.32)

    ax1 = fig.add_subplot(gs[0, 0])
    stages = ["Raw rows", "Qualified", "Balanced"]
    values = [
        int(dataset_manifest["raw_input_rows_total"]),
        int(dataset_manifest["num_samples_total"]) + int(dataset_manifest["phase_balancing"]["dropped_samples"]),
        int(dataset_manifest["num_samples_total"]),
    ]
    colors = ["#7aa6c2", "#8ec07c", "#e09f3e"]
    ypos = np.arange(len(stages))
    ax1.barh(ypos, values, color=colors, edgecolor="black")
    ax1.set_yticks(ypos, stages)
    ax1.invert_yaxis()
    ax1.set_title("Dataset funnel")
    ax1.set_xlabel("Patterns")
    for i, v in enumerate(values):
        ax1.text(v + 70, i, f"{v}", va="center", fontsize=16, weight="bold")

    ax2 = fig.add_subplot(gs[0, 1])
    split_names = ["Train", "Val", "Test"]
    phases = ["Al", "Cu", "Ni"]
    phase_colors = {"Al": "#4c78a8", "Cu": "#f58518", "Ni": "#54a24b"}
    x = np.arange(len(split_names))
    bottom = np.zeros(len(split_names))
    for phase in phases:
        vals = [int(dataset_manifest["split_phase_counts"][split.lower()][phase]) for split in split_names]
        ax2.bar(x, vals, bottom=bottom, color=phase_colors[phase], edgecolor="black", label=phase)
        bottom += np.array(vals)
    ax2.set_xticks(x, split_names)
    ax2.set_ylabel("Patterns")
    ax2.set_title("Balanced split composition")
    ax2.legend(loc="upper right", fontsize=14)
    for i, total in enumerate(bottom):
        ax2.text(i, total + 30, f"{int(total)}", ha="center", fontsize=16, weight="bold")

    ax3 = fig.add_subplot(gs[1, :])
    ax3.axis("off")
    rows = []
    for source in dataset_manifest["source_summaries"]:
        rows.append(
            [
                str(source["scan_id"]).replace("_", "-"),
                f"{int(source['rows_accepted'])} / {int(source['rows_total'])}",
                f"{float(source['accept_fraction']):.4f}",
                f"{tuple(source['pattern_shape'])}",
            ]
        )
    table = ax3.table(
        cellText=rows,
        colLabels=["Source scan", "Qualified / total", "Acceptance fraction", "Pattern shape"],
        loc="center",
        cellLoc="center",
        colLoc="center",
    )
    table.auto_set_font_size(False)
    table.set_fontsize(16)
    table.scale(1.1, 1.9)
    for key, cell in table.get_celld().items():
        cell.set_edgecolor("black")
        cell.set_linewidth(1.2)
        if key[0] == 0:
            cell.set_facecolor("#ededed")
            cell.set_text_props(weight="bold")
    ax3.set_title("Training source acceptance summary", fontsize=20, pad=12)

    return _savefig(fig, path)


def _figure_benchmark(suite_summary: dict[str, Any], path: Path) -> Path:
    rows = [row for row in suite_summary["rows"] if row.get("status") == "completed"]
    rows = sorted(rows, key=lambda row: float(row["test_macro_f1"]), reverse=True)
    names = [str(row["name"]) for row in rows]
    display_names = {
        "simple_cnn_w16": "cnn_w16",
        "simple_cnn_w32": "cnn_w32",
        "resnet18_scratch": "resnet18",
        "mobilenetv3_small_scratch": "mobilenetv3",
        "efficientnet_b0_scratch": "efficientnet_b0",
    }
    f1 = [float(row["test_macro_f1"]) for row in rows]
    runtime = [float(row["runtime_seconds"]) for row in rows]
    colors = ["#d62728" if name == "simple_cnn_w32" else "#4c78a8" for name in names]

    fig = plt.figure(figsize=(13.5, 7.2))
    gs = fig.add_gridspec(1, 2, width_ratios=[1.1, 1.0], wspace=0.28)

    ax1 = fig.add_subplot(gs[0, 0])
    ax1.scatter(runtime, f1, s=250, c=colors, edgecolors="black", linewidths=1.5)
    for x, y, name in zip(runtime, f1, names):
        ax1.annotate(display_names.get(name, name), (x, y), xytext=(6, 6), textcoords="offset points", fontsize=15)
    ax1.set_xlabel("Runtime (s)")
    ax1.set_ylabel("Test macro-F1")
    ax1.set_title("Accuracy vs runtime")
    ax1.set_xlim(min(runtime) - 5, max(runtime) + 20)
    ax1.set_ylim(0.95, 1.003)
    ax1.grid(alpha=0.25)

    ax2 = fig.add_subplot(gs[0, 1])
    ypos = np.arange(len(names))
    ax2.barh(ypos, f1, color=colors, edgecolor="black")
    ax2.set_yticks(ypos, [display_names.get(name, name) for name in names])
    ax2.invert_yaxis()
    ax2.set_xlim(0.95, 1.002)
    ax2.set_xlabel("Test macro-F1")
    ax2.set_title("Held-out ranking")
    for i, val in enumerate(f1):
        ax2.text(val + 0.0006, i, f"{val:.4f}", va="center", fontsize=15, weight="bold")

    return _savefig(fig, path)


def _draw_network(ax, *, x: float, y: float, w: float, h: float, title: str, lines: list[str], color: str) -> None:
    rect = patches.FancyBboxPatch(
        (x, y),
        w,
        h,
        boxstyle="round,pad=0.015,rounding_size=0.02",
        facecolor=color,
        edgecolor="black",
        linewidth=1.8,
    )
    ax.add_patch(rect)
    ax.text(x + 0.02, y + h - 0.06, title, fontsize=16, weight="bold", va="top")
    y0 = y + h - 0.13
    for idx, line in enumerate(lines):
        ax.text(x + 0.02, y0 - idx * 0.055, line, fontsize=13.5, va="top")


def _figure_model_architectures(path: Path) -> Path:
    fig, ax = plt.subplots(figsize=(13.5, 7.2))
    ax.set_axis_off()
    ax.text(0.03, 0.95, "Mermaid-style network overview", fontsize=24, weight="bold")
    ax.text(0.03, 0.90, "The gallery model is cnn_w32; the benchmark also compared larger timm backbones.", fontsize=18)

    # Top row
    _draw_network(
        ax,
        x=0.04,
        y=0.57,
        w=0.27,
        h=0.25,
        title="cnn_w16",
        lines=[
            "Input 128x128x1",
            "16-ch conv blocks",
            "Pooling -> 3-class head",
        ],
        color="#d9eaf7",
    )
    _draw_network(
        ax,
        x=0.365,
        y=0.57,
        w=0.27,
        h=0.25,
        title="cnn_w32",
        lines=[
            "Input 128x128x1",
            "Same layout, wider blocks",
            "Selected gallery model",
        ],
        color="#fce6c8",
    )
    _draw_network(
        ax,
        x=0.69,
        y=0.57,
        w=0.27,
        h=0.25,
        title="resnet18",
        lines=[
            "Input 128x128x1",
            "Residual blocks",
            "Top benchmark result",
        ],
        color="#e4f1d5",
    )

    # Bottom row
    _draw_network(
        ax,
        x=0.14,
        y=0.19,
        w=0.27,
        h=0.25,
        title="mobilenetv3_small",
        lines=[
            "Depthwise separable blocks",
            "Mobile-optimized design",
            "Near-top accuracy",
        ],
        color="#f4d6d6",
    )
    _draw_network(
        ax,
        x=0.52,
        y=0.19,
        w=0.27,
        h=0.25,
        title="efficientnet_b0",
        lines=[
            "Compound-scaled MBConv",
            "Largest suite runtime",
            "Top benchmark result",
        ],
        color="#d9dbf7",
    )

    connectors = [
        ((0.31, 0.695), (0.365, 0.695)),
        ((0.635, 0.695), (0.69, 0.695)),
        ((0.275, 0.44), (0.275, 0.57)),
        ((0.655, 0.44), (0.655, 0.57)),
    ]
    for start, end in connectors:
        ax.annotate("", xy=end, xytext=start, arrowprops=dict(arrowstyle="->", lw=2.0, color="black"))

    ax.text(0.035, 0.10, "Mermaid equivalent: input -> feature blocks -> global pooling -> 3-class head", fontsize=17, weight="bold")
    return _savefig(fig, path)


def _figure_selected_model(report: dict[str, Any], path: Path) -> Path:
    phases = ["Al", "Ni", "Cu"]
    cm = np.array(report["test_metrics"]["confusion_matrix"], dtype=float)
    cm_norm = cm / cm.sum(axis=1, keepdims=True)
    per_class = report["test_metrics"]["per_class"]
    metrics = np.array([[per_class[p]["precision"], per_class[p]["recall"], per_class[p]["f1"]] for p in phases], dtype=float)

    fig = plt.figure(figsize=(13.5, 7.2))
    gs = fig.add_gridspec(1, 2, width_ratios=[1.0, 1.05], wspace=0.3)

    ax1 = fig.add_subplot(gs[0, 0])
    im = ax1.imshow(cm_norm, cmap="Blues", vmin=0.0, vmax=1.0)
    ax1.set_xticks(range(len(phases)), phases)
    ax1.set_yticks(range(len(phases)), phases)
    ax1.set_xlabel("Predicted")
    ax1.set_ylabel("True")
    ax1.set_title("simple_cnn_w32 confusion matrix")
    for i in range(cm.shape[0]):
        for j in range(cm.shape[1]):
            ax1.text(
                j,
                i,
                f"{int(cm[i, j])}\n{cm_norm[i, j]:.3f}",
                ha="center",
                va="center",
                fontsize=15,
                color="white" if cm_norm[i, j] > 0.55 else "black",
                weight="bold",
            )
    fig.colorbar(im, ax=ax1, fraction=0.046, pad=0.04)

    ax2 = fig.add_subplot(gs[0, 1])
    x = np.arange(len(phases))
    width = 0.23
    labels = ["Precision", "Recall", "F1"]
    palette = ["#4c78a8", "#f58518", "#54a24b"]
    for idx in range(3):
        ax2.bar(x + (idx - 1) * width, metrics[:, idx], width=width, label=labels[idx], color=palette[idx], edgecolor="black")
    ax2.set_xticks(x, phases)
    ax2.set_ylim(0.97, 1.005)
    ax2.set_ylabel("Score")
    ax2.set_title("Per-class metrics")
    ax2.legend(loc="lower right", fontsize=14)
    for i in range(len(phases)):
        ax2.text(x[i], 0.9705, f"support={per_class[phases[i]]['support']}", ha="center", fontsize=14)

    return _savefig(fig, path)


def _figure_reference_gallery(gallery_root: Path, stats: dict[str, dict[str, Any]], path: Path) -> Path:
    contact = Image.open(gallery_root / "reference_contact_sheet.png")
    fig = plt.figure(figsize=(13.5, 7.2))
    gs = fig.add_gridspec(2, 1, height_ratios=[3.3, 1.15], hspace=0.15)

    ax1 = fig.add_subplot(gs[0, 0])
    ax1.imshow(contact)
    ax1.axis("off")
    ax1.set_title("Reference diagnostic gallery contact sheet")

    ax2 = fig.add_subplot(gs[1, 0])
    ax2.axis("off")
    rows = []
    for name in ["Al-2_1", "Cu-2_1", "Ni-2_1"]:
        row = stats[name]
        only_phase = list(row["counts"].keys())[0]
        rows.append(
            [
                name,
                f"5 / 5 {only_phase}",
                f"{row['mean_confidence']:.6f}",
                f"{row['mean_margin']:.6f}",
            ]
        )
    table = ax2.table(
        cellText=rows,
        colLabels=["Reference scan", "Predicted composition", "Mean confidence", "Mean margin"],
        loc="center",
        cellLoc="center",
        colLoc="center",
    )
    table.auto_set_font_size(False)
    table.set_fontsize(16)
    table.scale(1.0, 1.8)
    for key, cell in table.get_celld().items():
        cell.set_edgecolor("black")
        cell.set_linewidth(1.2)
        if key[0] == 0:
            cell.set_facecolor("#ededed")
            cell.set_text_props(weight="bold")

    return _savefig(fig, path)


def _figure_blind_gallery(gallery_root: Path, stats: dict[str, dict[str, Any]], path: Path) -> Path:
    contact = Image.open(gallery_root / "unknown_contact_sheet.png")
    fig = plt.figure(figsize=(13.5, 7.2))
    gs = fig.add_gridspec(2, 2, height_ratios=[2.6, 1.5], width_ratios=[1.2, 1.0], hspace=0.25, wspace=0.28)

    ax1 = fig.add_subplot(gs[0, :])
    ax1.imshow(contact)
    ax1.axis("off")
    ax1.set_title("Blind-scan diagnostic gallery contact sheet")

    ax2 = fig.add_subplot(gs[1, 0])
    scans = ["Data_1", "Data_2", "Data_3"]
    phase_order = ["Al", "Cu", "Ni"]
    colors = {"Al": "#4c78a8", "Cu": "#f58518", "Ni": "#54a24b"}
    x = np.arange(len(scans))
    bottom = np.zeros(len(scans))
    for phase in phase_order:
        vals = [stats[scan]["counts"].get(phase, 0) for scan in scans]
        ax2.bar(x, vals, bottom=bottom, color=colors[phase], edgecolor="black", label=phase)
        bottom += np.array(vals)
    ax2.set_xticks(x, scans)
    ax2.set_ylim(0, 5.6)
    ax2.set_ylabel("Votes in sampled tiles")
    ax2.set_title("Predicted composition of 5 sampled tiles")
    ax2.legend(loc="upper right", fontsize=13)
    for i, total in enumerate(bottom):
        ax2.text(i, total + 0.08, f"{int(total)}", ha="center", fontsize=15, weight="bold")

    ax3 = fig.add_subplot(gs[1, 1])
    ax3.axis("off")
    rows = []
    for scan in scans:
        rows.append(
            [
                scan,
                ", ".join(f"{phase}:{count}" for phase, count in sorted(stats[scan]["counts"].items())),
                f"{stats[scan]['mean_confidence']:.6f}",
                f"{stats[scan]['mean_margin']:.6f}",
            ]
        )
    table = ax3.table(
        cellText=rows,
        colLabels=["Blind scan", "Vote count", "Mean confidence", "Mean margin"],
        loc="center",
        cellLoc="center",
        colLoc="center",
    )
    table.auto_set_font_size(False)
    table.set_fontsize(15)
    table.scale(1.1, 1.85)
    for key, cell in table.get_celld().items():
        cell.set_edgecolor("black")
        cell.set_linewidth(1.2)
        if key[0] == 0:
            cell.set_facecolor("#ededed")
            cell.set_text_props(weight="bold")

    return _savefig(fig, path)


def _figure_takeaways(dataset_manifest: dict[str, Any], suite_summary: dict[str, Any], stats: dict[str, dict[str, Any]], path: Path) -> Path:
    fig, ax = plt.subplots(figsize=(13.5, 7.2))
    ax.set_axis_off()

    cards = [
        (0.04, 0.58, 0.42, 0.28, "#d9eaf7", "Balanced dataset", f"5253 selected patterns\n1751 per phase\n4200 / 525 / 528 split"),
        (0.54, 0.58, 0.42, 0.28, "#e4f1d5", "Benchmark leader", "resnet18_scratch and\nefficientnet_b0_scratch\n1.0000 test macro-F1"),
        (0.04, 0.16, 0.42, 0.28, "#fce6c8", "Deployment candidate", "simple_cnn_w32\n0.99621 test accuracy\n15.31 s runtime"),
        (0.54, 0.16, 0.42, 0.28, "#f4d6d6", "Blind-scan evidence", f"Data_3 -> 5 Al\nData_1 -> 4 Cu, 1 Ni\nData_2 -> 3 Ni, 2 Cu"),
    ]
    for x, y, w, h, color, title, body in cards:
        rect = patches.FancyBboxPatch((x, y), w, h, boxstyle="round,pad=0.015,rounding_size=0.03", facecolor=color, edgecolor="black", linewidth=2)
        ax.add_patch(rect)
        ax.text(x + 0.02, y + h - 0.07, title, fontsize=22, weight="bold", va="top")
        ax.text(x + 0.02, y + 0.07, body, fontsize=18, va="bottom")

    ax.text(0.04, 0.93, "Summary scorecard", fontsize=24, weight="bold")
    ax.text(0.04, 0.08, "Reference gallery agreement: 15 / 15 sampled tiles correct | Residual ambiguity remains on the Ni/Cu boundary.", fontsize=18)
    return _savefig(fig, path)


def _add_title_bar(slide, text: str) -> None:
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, Inches(13.333), Inches(0.8))
    shape.fill.solid()
    shape.fill.fore_color.rgb = BLACK
    shape.line.color.rgb = BLACK
    tf = shape.text_frame
    tf.clear()
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.name = "Arial"
    run.font.size = Pt(24)
    run.font.bold = True
    run.font.color.rgb = WHITE


def _add_takeaway_bar(slide, text: str) -> None:
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, Inches(6.8), Inches(13.333), Inches(0.7))
    shape.fill.solid()
    shape.fill.fore_color.rgb = BLACK
    shape.line.color.rgb = BLACK
    tf = shape.text_frame
    tf.clear()
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.name = "Arial"
    run.font.size = Pt(22)
    run.font.bold = True
    run.font.color.rgb = WHITE


def _add_bullets(slide, heading: str, bullets: list[str]) -> None:
    box = slide.shapes.add_textbox(Inches(0.35), Inches(1.0), Inches(2.45), Inches(5.55))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0

    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = heading
    run.font.name = "Arial"
    run.font.size = Pt(20)
    run.font.bold = True
    run.font.color.rgb = TEXT

    for bullet in bullets:
        p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.level = 0
        run = p.add_run()
        run.text = f"- {bullet}"
        run.font.name = "Arial"
        run.font.size = Pt(18)
        run.font.color.rgb = TEXT
        p.space_after = Pt(6)


def _add_image(slide, image_path: Path, *, left: float = 3.0, top: float = 1.0, width: float = 10.0, height: float = 5.6) -> None:
    slide.shapes.add_picture(str(image_path), Inches(left), Inches(top), width=Inches(width), height=Inches(height))


def _build_manifest(args: argparse.Namespace, outputs: dict[str, str]) -> dict[str, Any]:
    return {
        "deck_title": args.deck_title,
        "benchmark_root": str(args.benchmark_root.as_posix()),
        "gallery_root": str(args.gallery_root.as_posix()),
        "dataset_manifest": str(args.dataset_manifest.as_posix()),
        "artifacts": outputs,
        "style": {
            "title_font": {"name": "Arial", "size_pt": 24},
            "takeaway_font": {"name": "Arial", "size_pt": 22},
            "body_font": {"name": "Arial", "size_pt": 18},
            "subheading_font": {"name": "Arial", "size_pt": 20},
            "title_bar": "black",
            "takeaway_bar": "black",
        },
    }


def main() -> None:
    args = parse_args()
    _setup_matplotlib()

    repo_root = Path(__file__).resolve().parents[2]
    benchmark_root = args.benchmark_root if args.benchmark_root.is_absolute() else (repo_root / args.benchmark_root)
    gallery_root = args.gallery_root if args.gallery_root.is_absolute() else (repo_root / args.gallery_root)
    dataset_manifest_path = args.dataset_manifest if args.dataset_manifest.is_absolute() else (repo_root / args.dataset_manifest)
    output_dir = args.output_dir if args.output_dir.is_absolute() else (repo_root / args.output_dir)
    output_dir.mkdir(parents=True, exist_ok=True)

    suite_summary = _load_json(benchmark_root / "suite_summary.json")
    dataset_manifest = _load_json(dataset_manifest_path)
    gallery_manifest = _load_json(gallery_root / "manifest.json")
    simple_cnn_report = _load_json(benchmark_root / "simple_cnn_w32" / "report.json")
    gallery_stats = _gallery_stats(gallery_manifest)

    assets_dir = output_dir / f"{args.basename}_assets"
    assets_dir.mkdir(parents=True, exist_ok=True)

    outputs = {
        "workflow": str(_figure_workflow(assets_dir / "workflow_overview.png")),
        "dataset": str(_figure_dataset(dataset_manifest, assets_dir / "dataset_preparation.png")),
        "benchmark": str(_figure_benchmark(suite_summary, assets_dir / "benchmark_comparison.png")),
        "architectures": str(_figure_model_architectures(assets_dir / "model_architectures.png")),
        "selected_model": str(_figure_selected_model(simple_cnn_report, assets_dir / "selected_model_performance.png")),
        "reference_gallery": str(_figure_reference_gallery(gallery_root, gallery_stats, assets_dir / "reference_gallery.png")),
        "blind_gallery": str(_figure_blind_gallery(gallery_root, gallery_stats, assets_dir / "blind_gallery.png")),
        "takeaways": str(_figure_takeaways(dataset_manifest, suite_summary, gallery_stats, assets_dir / "takeaways.png")),
    }

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    slides = [
        {
            "title": "Workflow Overview",
            "heading": "Pipeline",
            "bullets": [
                "Three single-phase .oh5 scans for Al, Cu, and Ni were ingested.",
                "Quality filtering, phase balancing, and deterministic splitting were applied before training.",
                "Benchmark training was followed by diagnostic-gallery validation on reference and blind scans.",
            ],
            "image": outputs["workflow"],
            "takeaway": "The pipeline was fully reproducible from source scans to benchmark and blind-scan inspection.",
        },
        {
            "title": "Dataset Preparation",
            "heading": "Balanced build",
            "bullets": [
                "10025 raw rows became 7064 qualified rows after CI > 0.2 and Fit < 1.5.",
                "Balancing equalized the dataset to 1751 patterns per phase and dropped 1811 excess samples.",
                "The final split was 4200 train, 525 val, and 528 test with equal per-phase counts.",
            ],
            "image": outputs["dataset"],
            "takeaway": "The training benchmark used a deliberately balanced three-phase dataset, not a skewed class distribution.",
        },
        {
            "title": "Benchmark Model Results",
            "heading": "Five-model suite",
            "bullets": [
                "All five benchmark runs completed successfully on the same balanced split.",
                "resnet18_scratch and efficientnet_b0_scratch reached 1.0000 test macro-F1.",
                "simple_cnn_w32 reached 0.99621 test macro-F1 with much lower runtime than the larger backbones.",
            ],
            "image": outputs["benchmark"],
            "takeaway": "simple_cnn_w32 is the strongest lightweight model, while resnet18 and efficientnet_b0 set the benchmark ceiling.",
        },
        {
            "title": "Network Architectures",
            "heading": "Mermaid-style view",
            "bullets": [
                "cnn_w16 and cnn_w32 share the same compact topology, with cnn_w32 adding width rather than a new structure.",
                "resnet18, mobilenetv3_small, and efficientnet_b0 represent deeper timm backbones in the benchmark.",
                "cnn_w32 was chosen for gallery use because it balanced accuracy, compactness, and runtime.",
            ],
            "image": outputs["architectures"],
            "takeaway": "The selected gallery model kept a lightweight CNN structure while approaching the performance of much larger backbones.",
        },
        {
            "title": "simple_cnn_w32 Per-Class Behavior",
            "heading": "Selected model",
            "bullets": [
                "Held-out test accuracy was 0.99621, with perfect Al behavior and no Al confusion.",
                "The only residual errors were two Ni patterns misclassified as Cu.",
                "The remaining uncertainty is concentrated on the Ni/Cu boundary rather than the Al class.",
            ],
            "image": outputs["selected_model"],
            "takeaway": "The selected lightweight model is highly accurate, with its remaining weakness localized to Ni versus Cu separation.",
        },
        {
            "title": "Diagnostic Gallery Verification: Reference",
            "heading": "Reference scans",
            "bullets": [
                "The gallery used simple_cnn_w32 with CI > 0.5, Fit < 1.5, confidence > 0.5, and margin > 0.15.",
                "Five sampled tiles were drawn from each of Al-2_1, Cu-2_1, and Ni-2_1 using seed 7.",
                "All 15 sampled reference tiles matched the expected phase identity with very high confidence and margin.",
            ],
            "image": outputs["reference_gallery"],
            "takeaway": "The gallery model stayed internally consistent on the known reference scans under the stricter gallery filters.",
        },
        {
            "title": "Diagnostic Gallery Verification: Blind",
            "heading": "Blind scans",
            "bullets": [
                "Data_3 was unanimously predicted as Al across all five sampled tiles with high confidence and margin.",
                "Data_1 leaned Cu with one Ni vote, while Data_2 split between Ni and Cu with Ni the majority.",
                "These outputs are sampled-pattern evidence only and should not be treated as whole-scan proof.",
            ],
            "image": outputs["blind_gallery"],
            "takeaway": "Blind-scan evidence is encouraging, but only Data_3 is cleanly resolved in the current sampled gallery output.",
        },
        {
            "title": "Major Outcomes and Limitations",
            "heading": "Bottom line",
            "bullets": [
                "Balanced-data training substantially improved the ML branch and produced a strong five-model suite.",
                "simple_cnn_w32 is the best model to highlight for lightweight deployment and gallery use.",
                "Further blind-scan confidence still depends on reducing mixed Ni/Cu behavior and moving beyond five sampled tiles.",
            ],
            "image": outputs["takeaways"],
            "takeaway": "The ML workflow is now strong enough for lab presentation, with clear progress and clearly bounded remaining risk.",
        },
    ]

    for spec in slides:
        slide = prs.slides.add_slide(blank)
        _add_title_bar(slide, spec["title"])
        _add_bullets(slide, spec["heading"], spec["bullets"])
        _add_image(slide, Path(spec["image"]))
        _add_takeaway_bar(slide, spec["takeaway"])

    ppt_path = output_dir / f"{args.basename}.pptx"
    prs.save(ppt_path)

    manifest = _build_manifest(args, outputs)
    manifest_path = output_dir / f"{args.basename}_manifest.json"
    manifest_path.write_text(json.dumps(manifest, indent=2), encoding="utf-8")

    print(f"[ok] pptx written: {ppt_path}")
    print(f"[ok] manifest written: {manifest_path}")
    print(f"[ok] assets dir: {assets_dir}")


if __name__ == "__main__":
    main()
