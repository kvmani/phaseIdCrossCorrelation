#!/usr/bin/env python3
"""Generate a concise lab-meeting PPTX from ML benchmark artifacts."""

from __future__ import annotations

import argparse
import json
from pathlib import Path
import re
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description="Generate a lab-meeting PPT from ML benchmark artifacts.")
    parser.add_argument("--scan-root", required=True, help="Benchmark suite output root.")
    parser.add_argument("--output-dir", required=True, help="Directory for manifest and PPTX.")
    parser.add_argument("--deck-title", default=None, help="Deck title.")
    parser.add_argument("--basename", default=None, help="Output basename without extension.")
    parser.add_argument("--max-results", type=int, default=8, help="Unused compatibility argument.")
    return parser.parse_args()


def _slugify(text: str) -> str:
    text = re.sub(r"[^a-zA-Z0-9]+", "-", text.strip().lower())
    return re.sub(r"-+", "-", text).strip("-") or "ml-results"


def _load_json(path: Path) -> dict[str, Any]:
    return json.loads(path.read_text(encoding="utf-8"))


def _add_bar(slide, *, top: float, height: float, text: str, font_size: int) -> None:
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(top), Inches(13.333), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0, 0, 0)
    shape.line.color.rgb = RGBColor(0, 0, 0)
    tf = shape.text_frame
    tf.clear()
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.2)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    p.text = text
    run = p.runs[0]
    run.font.name = "Arial"
    run.font.size = Pt(font_size)
    run.font.bold = True
    run.font.color.rgb = RGBColor(255, 255, 255)


def _add_bullets(slide, bullets: list[str]) -> None:
    box = slide.shapes.add_textbox(Inches(0.35), Inches(1.0), Inches(2.35), Inches(5.55))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    for idx, bullet in enumerate(bullets[:6]):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = f"• {bullet}"
        p.alignment = PP_ALIGN.LEFT
        run = p.runs[0]
        run.font.name = "Arial"
        run.font.size = Pt(18)


def _add_body_text(slide, title: str, lines: list[str], bottom_line: str) -> None:
    _add_bar(slide, top=0.0, height=0.8, text=title, font_size=24)
    _add_bullets(slide, lines)

    body = slide.shapes.add_textbox(Inches(2.95), Inches(1.05), Inches(10.0), Inches(5.4))
    tf = body.text_frame
    tf.clear()
    tf.word_wrap = True
    for idx, line in enumerate(lines):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = line
        p.alignment = PP_ALIGN.LEFT
        run = p.runs[0]
        run.font.name = "Arial"
        run.font.size = Pt(18)

    _add_bar(slide, top=6.72, height=0.7, text=bottom_line, font_size=20)


def _build_manifest(deck_title: str, suite_summary: dict[str, Any]) -> dict[str, Any]:
    rows = suite_summary.get("rows", [])
    completed = [row for row in rows if row.get("status") == "completed"]
    ranked = sorted(completed, key=lambda row: float(row.get("best_val_macro_f1") or -1.0), reverse=True)
    best = ranked[0] if ranked else None

    dataset_runs = suite_summary.get("runs_total", 0)
    slides: list[dict[str, Any]] = [
        {
            "title": deck_title,
            "lines": [
                "Objective: compare candidate classifiers on the prepared Al/Ni/Cu EBSD dataset.",
                f"Runs completed: {suite_summary.get('runs_completed', 0)} / {dataset_runs}",
                f"Runs failed: {suite_summary.get('runs_failed', 0)}",
                f"Best run: {suite_summary.get('best_run')}",
            ],
            "bottom_line": "Suite output gives a direct model-to-model comparison for lab review.",
        },
        {
            "title": "Methodology Overview",
            "lines": [
                "Input: filtered .oh5 patterns after CI/Fit quality gating.",
                "Preprocessing: resize, circular mask, per-pattern normalization.",
                "Evaluation: validation macro-F1 plus held-out test macro-F1.",
                "Outputs: per-run report.json, checkpoints, suite_summary.json, HTML report.",
            ],
            "bottom_line": "The workflow keeps machine-readable and human-readable artifacts aligned.",
        },
    ]

    comparison_lines = []
    for row in ranked[:6]:
        comparison_lines.append(
            f"{row.get('name')}: val_macro_f1={float(row.get('best_val_macro_f1') or 0.0):.4f}, "
            f"test_macro_f1={float(row.get('test_macro_f1') or 0.0):.4f}, "
            f"runtime={float(row.get('runtime_seconds') or 0.0):.2f}s"
        )
    slides.append(
        {
            "title": "Model Comparison",
            "lines": comparison_lines or ["No completed runs found."],
            "bottom_line": "Use validation macro-F1 and runtime together when selecting the next production model.",
        }
    )

    if best:
        slides.append(
            {
                "title": "Best Run",
                "lines": [
                    f"Run name: {best.get('name')}",
                    f"Model: {best.get('model_name')}",
                    f"Best validation macro-F1: {float(best.get('best_val_macro_f1') or 0.0):.4f}",
                    f"Test accuracy: {float(best.get('test_accuracy') or 0.0):.4f}",
                    f"Test macro-F1: {float(best.get('test_macro_f1') or 0.0):.4f}",
                    f"Runtime: {float(best.get('runtime_seconds') or 0.0):.2f}s",
                ],
                "bottom_line": "This run is the current benchmark leader on the prepared dataset.",
            }
        )

    slides.append(
        {
            "title": "Next Steps",
            "lines": [
                "Replace placeholder OH5 file names with the final production scan names.",
                "Increase epochs and, if needed, hold out 100 samples per phase for val/test.",
                "Rerun the benchmark suite and regenerate the deck from the updated suite output.",
            ],
            "bottom_line": "Only small config edits should be needed before the full production run.",
        }
    )

    return {
        "deck_title": deck_title,
        "created_utc": None,
        "slides": slides,
        "suite_summary": suite_summary,
    }


def _write_ppt(manifest: dict[str, Any], output_path: Path) -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    layout = prs.slide_layouts[6]

    for slide_spec in manifest["slides"]:
        slide = prs.slides.add_slide(layout)
        _add_body_text(
            slide,
            title=str(slide_spec["title"]),
            lines=[str(x) for x in slide_spec["lines"]],
            bottom_line=str(slide_spec["bottom_line"]),
        )

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_path)


def main() -> None:
    args = parse_args()
    scan_root = Path(args.scan_root).expanduser().resolve()
    output_dir = Path(args.output_dir).expanduser().resolve()
    output_dir.mkdir(parents=True, exist_ok=True)

    suite_summary_path = scan_root / "suite_summary.json"
    if not suite_summary_path.exists():
        raise SystemExit(f"suite_summary.json not found under {scan_root}")

    suite_summary = _load_json(suite_summary_path)
    deck_title = args.deck_title or f"{scan_root.name} Lab Meeting"
    base = args.basename or _slugify(deck_title)

    manifest = _build_manifest(deck_title, suite_summary)
    manifest_path = output_dir / f"{base}_manifest.json"
    ppt_path = output_dir / f"{base}.pptx"
    manifest_path.write_text(json.dumps(manifest, indent=2), encoding="utf-8")
    _write_ppt(manifest, ppt_path)

    print(f"[ok] manifest written: {manifest_path}")
    print(f"[ok] total slides: {len(manifest['slides'])}")
    print(f"[ok] pptx written: {ppt_path}")


if __name__ == "__main__":
    main()
